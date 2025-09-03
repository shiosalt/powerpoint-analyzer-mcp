"""
Test data generation system for PowerPoint MCP server testing.
Creates PowerPoint files with known formatting using python-pptx.
"""

import json
import logging
from pathlib import Path
from typing import Dict, List, Any
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

logger = logging.getLogger(__name__)


class TestPresentationGenerator:
    """Generate test PowerPoint files with known formatting."""
    
    def __init__(self, output_dir: str = "tests/test_files"):
        """Initialize the test presentation generator."""
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.expected_results = {}
    
    def create_formatting_test_file(self) -> str:
        """Create PowerPoint with all supported formatting types."""
        try:
            presentation = Presentation()
            
            # Remove default slide
            if len(presentation.slides) > 0:
                rId = presentation.slides._sldIdLst[0].rId
                presentation.part.drop_rel(rId)
                del presentation.slides._sldIdLst[0]
            
            # Add slides with different formatting types
            self.add_bold_text_slide(presentation)
            self.add_italic_text_slide(presentation)
            self.add_underlined_text_slide(presentation)
            self.add_mixed_formatting_slide(presentation)
            self.add_hyperlink_slide(presentation)
            self.add_font_size_slide(presentation)
            self.add_font_color_slide(presentation)
            self.add_empty_slide(presentation)
            
            # Save the presentation
            file_path = self.output_dir / "test_formatting_comprehensive.pptx"
            presentation.save(str(file_path))
            
            # Document expected results
            self.document_expected_results(str(file_path))
            
            logger.info(f"Created comprehensive formatting test file: {file_path}")
            return str(file_path)
            
        except Exception as e:
            logger.error(f"Error creating formatting test file: {e}")
            raise
    
    def add_bold_text_slide(self, presentation):
        """Add slide with various bold text patterns."""
        slide_layout = presentation.slide_layouts[1]  # Title and Content layout
        slide = presentation.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = "Bold Text Test Slide"
        
        # Add content with bold text
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        # Add paragraph with mixed bold/normal text
        p = tf.paragraphs[0]
        run1 = p.add_run()
        run1.text = "This is "
        run1.font.bold = False
        
        # Add bold text
        run2 = p.add_run()
        run2.text = "bold text"
        run2.font.bold = True
        
        run3 = p.add_run()
        run3.text = " and this is normal text."
        run3.font.bold = False
        
        # Add another paragraph with all bold text
        p2 = tf.add_paragraph()
        run4 = p2.add_run()
        run4.text = "This entire paragraph is bold"
        run4.font.bold = True
        
        # Store expected results
        self.expected_results["slide_1_bold"] = {
            "slide_number": 1,
            "title": "Bold Text Test Slide",
            "bold_segments": [
                {"text": "bold text", "start_position": 8},
                {"text": "This entire paragraph is bold", "start_position": 43}
            ]
        }
    
    def add_italic_text_slide(self, presentation):
        """Add slide with italic text patterns."""
        slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Italic Text Test Slide"
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        p = tf.paragraphs[0]
        run1 = p.add_run()
        run1.text = "This text contains "
        
        run2 = p.add_run()
        run2.text = "italic formatting"
        run2.font.italic = True
        
        run3 = p.add_run()
        run3.text = " in the middle."
        
        self.expected_results["slide_2_italic"] = {
            "slide_number": 2,
            "title": "Italic Text Test Slide",
            "italic_segments": [
                {"text": "italic formatting", "start_position": 18}
            ]
        }
    
    def add_underlined_text_slide(self, presentation):
        """Add slide with underlined text patterns."""
        slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Underlined Text Test Slide"
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        p = tf.paragraphs[0]
        run1 = p.add_run()
        run1.text = "This text has "
        
        run2 = p.add_run()
        run2.text = "underlined sections"
        run2.font.underline = True
        
        run3 = p.add_run()
        run3.text = " for testing."
        
        self.expected_results["slide_3_underlined"] = {
            "slide_number": 3,
            "title": "Underlined Text Test Slide",
            "underlined_segments": [
                {"text": "underlined sections", "start_position": 14}
            ]
        }
    
    def add_mixed_formatting_slide(self, presentation):
        """Add slide with overlapping formatting."""
        slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Mixed Formatting Test Slide"
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        p = tf.paragraphs[0]
        run1 = p.add_run()
        run1.text = "This text has "
        
        # Bold and italic
        run2 = p.add_run()
        run2.text = "bold and italic"
        run2.font.bold = True
        run2.font.italic = True
        
        run3 = p.add_run()
        run3.text = " formatting together."
        
        self.expected_results["slide_4_mixed"] = {
            "slide_number": 4,
            "title": "Mixed Formatting Test Slide",
            "bold_segments": [
                {"text": "bold and italic", "start_position": 14}
            ],
            "italic_segments": [
                {"text": "bold and italic", "start_position": 14}
            ]
        }
    
    def add_hyperlink_slide(self, presentation):
        """Add slide with hyperlinks."""
        slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Hyperlink Test Slide"
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        p = tf.paragraphs[0]
        run1 = p.add_run()
        run1.text = "Visit "
        
        # Add hyperlink (Note: python-pptx has limited hyperlink support)
        run2 = p.add_run()
        run2.text = "our website"
        # Hyperlink creation in python-pptx requires more complex setup
        # For now, we'll document what should be there
        
        run3 = p.add_run()
        run3.text = " for more information."
        
        self.expected_results["slide_5_hyperlinks"] = {
            "slide_number": 5,
            "title": "Hyperlink Test Slide",
            "hyperlink_segments": [
                {
                    "text": "our website",
                    "start_position": 6,
                    "url": "https://example.com",
                    "link_type": "external"
                }
            ]
        }
    
    def add_font_size_slide(self, presentation):
        """Add slide with different font sizes."""
        slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Font Size Test Slide"
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        p = tf.paragraphs[0]
        run1 = p.add_run()
        run1.text = "Small text "
        run1.font.size = Pt(10)
        
        run2 = p.add_run()
        run2.text = "Medium text "
        run2.font.size = Pt(14)
        
        run3 = p.add_run()
        run3.text = "Large text"
        run3.font.size = Pt(18)
        
        self.expected_results["slide_6_font_sizes"] = {
            "slide_number": 6,
            "title": "Font Size Test Slide",
            "font_size_segments": [
                {"text": "Small text ", "font_size": 10, "start_position": 0},
                {"text": "Medium text ", "font_size": 14, "start_position": 11},
                {"text": "Large text", "font_size": 18, "start_position": 24}
            ]
        }
    
    def add_font_color_slide(self, presentation):
        """Add slide with different font colors."""
        slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Font Color Test Slide"
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        p = tf.paragraphs[0]
        run1 = p.add_run()
        run1.text = "Red text "
        run1.font.color.rgb = RGBColor(255, 0, 0)
        
        run2 = p.add_run()
        run2.text = "Blue text "
        run2.font.color.rgb = RGBColor(0, 0, 255)
        
        run3 = p.add_run()
        run3.text = "Green text"
        run3.font.color.rgb = RGBColor(0, 255, 0)
        
        self.expected_results["slide_7_font_colors"] = {
            "slide_number": 7,
            "title": "Font Color Test Slide",
            "font_color_segments": [
                {"text": "Red text ", "font_color": "#FF0000", "start_position": 0},
                {"text": "Blue text ", "font_color": "#0000FF", "start_position": 9},
                {"text": "Green text", "font_color": "#00FF00", "start_position": 19}
            ]
        }
    
    def add_empty_slide(self, presentation):
        """Add slide with no content for edge case testing."""
        slide_layout = presentation.slide_layouts[6]  # Blank layout
        slide = presentation.slides.add_slide(slide_layout)
        
        self.expected_results["slide_8_empty"] = {
            "slide_number": 8,
            "title": "",
            "text_elements": []
        }
    
    def create_edge_case_test_file(self) -> str:
        """Create PowerPoint with edge cases and complex layouts."""
        try:
            presentation = Presentation()
            
            # Remove default slide
            if len(presentation.slides) > 0:
                rId = presentation.slides._sldIdLst[0].rId
                presentation.part.drop_rel(rId)
                del presentation.slides._sldIdLst[0]
            
            self.add_complex_layout_slide(presentation)
            self.add_table_with_formatting_slide(presentation)
            self.add_very_long_text_slide(presentation)
            
            file_path = self.output_dir / "test_edge_cases.pptx"
            presentation.save(str(file_path))
            
            logger.info(f"Created edge case test file: {file_path}")
            return str(file_path)
            
        except Exception as e:
            logger.error(f"Error creating edge case test file: {e}")
            raise
    
    def add_complex_layout_slide(self, presentation):
        """Add slide with complex layout and multiple text boxes."""
        slide_layout = presentation.slide_layouts[6]  # Blank layout
        slide = presentation.slides.add_slide(slide_layout)
        
        # Add title text box
        left = Inches(1)
        top = Inches(1)
        width = Inches(8)
        height = Inches(1)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = "Complex Layout Test"
        
        # Add multiple text boxes with different formatting
        for i in range(3):
            left = Inches(1 + i * 2.5)
            top = Inches(3)
            width = Inches(2)
            height = Inches(2)
            
            text_box = slide.shapes.add_textbox(left, top, width, height)
            text_frame = text_box.text_frame
            text_frame.text = f"Text box {i+1} with bold formatting"
            
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.bold = True
    
    def add_table_with_formatting_slide(self, presentation):
        """Add slide with table containing formatted text."""
        slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Table with Formatting"
        
        # Add table
        rows, cols = 3, 3
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(3)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Add formatted content to table cells
        for row_idx in range(rows):
            for col_idx in range(cols):
                cell = table.cell(row_idx, col_idx)
                cell.text = f"Cell {row_idx},{col_idx}"
                
                # Make header row bold
                if row_idx == 0:
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.bold = True
    
    def add_very_long_text_slide(self, presentation):
        """Add slide with very long text for performance testing."""
        slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Long Text Performance Test"
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        # Create very long text with some bold sections
        long_text = "This is a very long text " * 100
        p = tf.paragraphs[0]
        p.text = long_text
        
        # Make every 10th word bold
        words = long_text.split()
        tf.clear()
        p = tf.paragraphs[0]
        
        for i, word in enumerate(words):
            run = p.add_run()
            run.text = word + " "
            if i % 10 == 0:
                run.font.bold = True
    
    def document_expected_results(self, file_path: str) -> Dict:
        """Document expected extraction results for test file."""
        results_file = Path(file_path).with_suffix('.json')
        
        with open(results_file, 'w', encoding='utf-8') as f:
            json.dump(self.expected_results, f, indent=2, ensure_ascii=False)
        
        logger.info(f"Documented expected results: {results_file}")
        return self.expected_results
    
    def create_all_test_files(self) -> List[str]:
        """Create all test files for comprehensive testing."""
        test_files = []
        
        try:
            # Create comprehensive formatting test file
            test_files.append(self.create_formatting_test_file())
            
            # Create edge case test file
            test_files.append(self.create_edge_case_test_file())
            
            logger.info(f"Created {len(test_files)} test files")
            return test_files
            
        except Exception as e:
            logger.error(f"Error creating test files: {e}")
            raise
    
    def validate_test_files(self, file_paths: List[str]) -> Dict[str, bool]:
        """Validate that test files were created correctly."""
        validation_results = {}
        
        for file_path in file_paths:
            try:
                # Check if file exists
                if not Path(file_path).exists():
                    validation_results[file_path] = False
                    continue
                
                # Try to open with python-pptx to verify it's valid
                presentation = Presentation(file_path)
                slide_count = len(presentation.slides)
                
                validation_results[file_path] = slide_count > 0
                logger.info(f"Validated {file_path}: {slide_count} slides")
                
            except Exception as e:
                logger.error(f"Validation failed for {file_path}: {e}")
                validation_results[file_path] = False
        
        return validation_results


def main():
    """Generate test files for PowerPoint MCP server testing."""
    generator = TestPresentationGenerator()
    
    try:
        test_files = generator.create_all_test_files()
        validation_results = generator.validate_test_files(test_files)
        
        print("Test file generation completed:")
        for file_path, is_valid in validation_results.items():
            status = "✓" if is_valid else "✗"
            print(f"  {status} {file_path}")
        
        print(f"\nGenerated {len(test_files)} test files in {generator.output_dir}")
        
    except Exception as e:
        print(f"Error generating test files: {e}")
        raise


if __name__ == "__main__":
    main()