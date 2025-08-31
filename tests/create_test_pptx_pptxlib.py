"""
python-pptx を使って、テスト用 PowerPoint ファイルを作成するスクリプト。
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pathlib import Path

def create_minimal_pptx(output_path: str):
    """最小限のスライドを持つ PPTX を作成"""
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]  # タイトルスライド
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Test Presentation Title"
    slide.placeholders[1].text = "Integration Test Subtitle"
    prs.save(output_path)

def create_complex_pptx(output_path: str):
    """テーブルと複数スライドを持つ PPTX を作成"""
    prs = Presentation()
    # 1枚目: タイトルスライド
    slide_layout = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(slide_layout)
    slide1.shapes.title.text = "Complex Test Presentation"
    # 2枚目: テーブル付きスライド
    slide_layout2 = prs.slide_layouts[5]  # タイトルとコンテンツ
    slide2 = prs.slides.add_slide(slide_layout2)
    shapes = slide2.shapes
    title_shape = shapes.title
    title_shape.text = "Table Slide"
    rows, cols = 3, 2
    left = Inches(2)
    top = Inches(2)
    width = Inches(4)
    height = Inches(1)
    table = shapes.add_table(rows, cols, left, top, width, height).table
    # ヘッダー
    table.cell(0, 0).text = "Header 1"
    table.cell(0, 1).text = "Header 2"
    # データ
    table.cell(1, 0).text = "Data 1"
    table.cell(1, 1).text = "Data 2"
    table.cell(2, 0).text = "Data 3"
    table.cell(2, 1).text = "Data 4"
    prs.save(output_path)

def create_large_pptx(output_path: str):
    """パフォーマンステスト用の大きな PPTX を作成"""
    prs = Presentation()
    for i in range(20):
        slide_layout = prs.slide_layouts[1]  # タイトルとコンテンツ
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = f"Large Test Presentation - Slide {i+1}"
        content = f"This is slide {i+1} with lots of content. " * 100
        slide.placeholders[1].text = content
    prs.save(output_path)

if __name__ == "__main__":
    test_dir = Path("tests/test_files")
    test_dir.mkdir(exist_ok=True)
    create_minimal_pptx(str(test_dir / "test_minimal.pptx"))
    create_complex_pptx(str(test_dir / "test_complex.pptx"))
    create_large_pptx(str(test_dir / "test_large.pptx"))
    print("Created test PowerPoint files:")
    print(f"- {test_dir / 'test_minimal.pptx'}")
    print(f"- {test_dir / 'test_complex.pptx'}")
    print(f"- {test_dir / 'test_large.pptx'}")
